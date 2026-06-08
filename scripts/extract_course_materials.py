#!/usr/bin/env python3
from __future__ import annotations

import argparse
import hashlib
import html
import json
import re
import time
from pathlib import Path

from docx import Document
from pdfminer.high_level import extract_text as extract_pdf_text
from pptx import Presentation


DEFAULT_ROOT = Path("/Users/mac/Desktop/教学与课程材料/讲义（講義）/日本語翻訳　講義資料")
MAX_TEXT = 7000

EXCLUDED_PARTS = {
    "二班宿題",
    "历年学生成绩",
    "期末試験",
    "期末考试",
    "试题",
    "作业",
    "实践调研活动",
    "他山之石",
}

EXCLUDED_NAME_PATTERNS = [
    r"^~\$",
    r"^\.~",
    r"成绩",
    r"名单",
    r"分组",
    r"登记表",
    r"试卷",
    r"试题",
    r"答案",
    r"审批",
    r"学生",
    r"宿題",
    r"記入シート",
    r"参考",
    r"バックアップ",
    r"备份",
    r"风险",
    r"责任",
    r"计划书",
]

ALLOWED_SUFFIXES = {".pptx", ".docx", ".pdf", ".html"}


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--root", type=Path, default=DEFAULT_ROOT)
    parser.add_argument("--output", type=Path, default=Path("data/course_materials.json"))
    parser.add_argument("--limit", type=int, default=80)
    args = parser.parse_args()

    files = select_files(args.root)[: args.limit]
    materials = []
    for path in files:
        try:
            text, meta = extract_material(path)
        except Exception as exc:
            text, meta = "", {"error": str(exc)}
        if len(text.strip()) < 40:
            continue
        materials.append(build_material(path, text, meta, args.root))
    materials = dedupe_materials(materials)

    args.output.parent.mkdir(parents=True, exist_ok=True)
    args.output.write_text(json.dumps({"replace": True, "materials": materials}, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"wrote {len(materials)} materials to {args.output}")


def select_files(root: Path) -> list[Path]:
    candidates = []
    for path in root.rglob("*"):
        if not path.is_file() or path.suffix.lower() not in ALLOWED_SUFFIXES:
            continue
        if should_exclude(path, root):
            continue
        if not is_course_material(path, root):
            continue
        candidates.append(path)
    return sorted(candidates, key=sort_key)


def should_exclude(path: Path, root: Path) -> bool:
    relative_parts = set(path.relative_to(root).parts)
    if relative_parts & EXCLUDED_PARTS:
        return True
    name = path.name
    return any(re.search(pattern, name, re.I) for pattern in EXCLUDED_NAME_PATTERNS)


def is_course_material(path: Path, root: Path) -> bool:
    rel = str(path.relative_to(root))
    if "日语笔译第一学期" in rel:
        return bool(re.search(r"第[一二三四五六七八九十十一十二十三十四0-9]+[課课]|尹国鵬 20年日语笔译课程|可能态|形式名词|惯用|成语|被动|长句|拟声", rel))
    if "第二学期/2026" in rel:
        return True
    if "第二学期/笔译2 课件汇总" in rel:
        return True
    if "第二学期" in rel and path.suffix.lower() == ".pptx":
        return bool(re.search(r"概論|致辞|商务|広告|广告|ゲーム|新聞|技術|旅游|観光|法律|招商|入札", path.name))
    return False


def sort_key(path: Path):
    week = infer_week(path)
    return (week or 99, canonical_topic(path), -path.stat().st_mtime, str(path))


def extract_material(path: Path) -> tuple[str, dict]:
    suffix = path.suffix.lower()
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".pdf":
        text = extract_pdf_text(str(path), maxpages=8)
        return clean_text(text), {"pages_sampled": 8}
    if suffix == ".html":
        return extract_html(path), {}
    return "", {}


def extract_pptx(path: Path) -> tuple[str, dict]:
    prs = Presentation(str(path))
    chunks = []
    for slide_no, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        if texts:
            chunks.append(f"Slide {slide_no}\n" + "\n".join(texts))
    return clean_text("\n\n".join(chunks)), {"slides": len(prs.slides)}


def extract_docx(path: Path) -> tuple[str, dict]:
    doc = Document(str(path))
    paragraphs = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
    table_text = []
    for table in doc.tables[:8]:
        for row in table.rows[:20]:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                table_text.append(" | ".join(cells))
    return clean_text("\n".join(paragraphs + table_text)), {"paragraphs": len(paragraphs), "tables_sampled": min(len(doc.tables), 8)}


def extract_html(path: Path) -> str:
    raw = path.read_text(encoding="utf-8", errors="ignore")
    raw = re.sub(r"<(script|style).*?</\1>", " ", raw, flags=re.I | re.S)
    raw = re.sub(r"<[^>]+>", " ", raw)
    return clean_text(html.unescape(raw))


def clean_text(value: str) -> str:
    value = re.sub(r"\r", "\n", value)
    value = redact_private_contact(value)
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()[:MAX_TEXT]


def redact_private_contact(value: str) -> str:
    value = re.sub(r"1[3-9]\d{9}", "[已脱敏手机号]", value)
    value = re.sub(r"[\w.+-]+[@＠][\w.-]+", "[已脱敏邮箱]", value)
    value = re.sub(r"(?im)^.*(?:wechat|微信|連絡先|联系方式).*$", "[已脱敏联系方式]", value)
    return value


def build_material(path: Path, text: str, meta: dict, root: Path) -> dict:
    rel = str(path.relative_to(root))
    week = infer_week(path)
    title = title_from_path(path)
    category = category_from_path(path)
    digest = hashlib.sha1(str(path).encode("utf-8")).hexdigest()[:12]
    modified_at = path.stat().st_mtime
    return {
        "title": title,
        "category": category,
        "material_type": path.suffix.lower().lstrip("."),
        "week": week,
        "source_path": rel,
        "summary": summarize(text, title),
        "content": text,
        "metadata": {
            "sha": digest,
            "file_name": path.name,
            "source_root": str(root),
            "modified_at": int(modified_at),
            "modified_at_text": time.strftime("%Y-%m-%d %H:%M:%S", time.localtime(modified_at)),
            "canonical_topic": canonical_topic(path),
            **meta,
        },
    }


def dedupe_materials(materials: list[dict]) -> list[dict]:
    by_key: dict[tuple, dict] = {}
    for item in materials:
        key = material_dedupe_key(item)
        current = by_key.get(key)
        if current is None or material_rank(item) > material_rank(current):
            by_key[key] = item
    return sorted(by_key.values(), key=lambda item: (item.get("week") or 99, item["category"], item["title"]))


def material_dedupe_key(item: dict) -> tuple:
    topic = item["metadata"].get("canonical_topic")
    week = 1 if topic == "intro" else item.get("week")
    return (week, item["category"], topic)


def material_rank(item: dict) -> tuple:
    meta = item.get("metadata", {})
    type_score = {"html": 4, "pptx": 3, "docx": 2, "pdf": 1}.get(item.get("material_type"), 0)
    return (int(meta.get("modified_at") or 0), type_score, len(item.get("content") or ""))


def title_from_path(path: Path) -> str:
    return re.sub(r"[_　 ]+", " ", path.stem).strip()


def canonical_topic(path: Path) -> str:
    text = str(path)
    name = path.stem.lower()
    topic_rules = [
        ("intro", r"序論|绪论|概论|概論|第一課|第一课"),
        ("ja-zh-diff", r"日汉主要不同|主要不同点"),
        ("kanji", r"汉字词|漢字詞"),
        ("addition-omission", r"增词|减词|引申"),
        ("conversion", r"转换词语|正反表达|やいなや"),
        ("onomatopoeia", r"拟声|拟态|オノマト"),
        ("loanwords", r"外来语"),
        ("pronouns", r"人称代词|指示词|一人称|二人称"),
        ("formal-nouns", r"形式名词|形式名詞"),
        ("idioms-proverbs", r"成语|谚语|ことわざ|第十[課课]"),
        ("idioms", r"习惯语|惯用"),
        ("causative", r"使役"),
        ("passive", r"被动|passive|pasive"),
        ("long-sentence", r"长句"),
        ("speech", r"致辞"),
        ("game-localization", r"游戏|ゲーム|ローカライズ"),
        ("advertising", r"广告|広告"),
        ("business-doc", r"商务|ビジネス|信函|文書"),
        ("legal", r"法律|法規"),
        ("tourism", r"旅游|観光"),
        ("news", r"新聞|新闻"),
        ("technical", r"技術|技术"),
        ("bidding", r"招商|招标|招投标|入札"),
    ]
    for topic, pattern in topic_rules:
        if re.search(pattern, text, re.I) or re.search(pattern, name, re.I):
            return topic
    return re.sub(r"[^a-z0-9一-龥ぁ-んァ-ン]+", "-", name).strip("-")


def category_from_path(path: Path) -> str:
    text = str(path)
    if "日语笔译第一学期" in text:
        return "日语笔译第一学期"
    if "第二学期" in text:
        return "日语笔译第二学期"
    if "日汉翻译教程高宁" in text:
        return "日汉翻译教程参考"
    return "课程资料"


def infer_week(path: Path) -> int | None:
    text = str(path)
    cn_map = {
        "一": 1,
        "二": 2,
        "三": 3,
        "四": 4,
        "五": 5,
        "六": 6,
        "七": 7,
        "八": 8,
        "九": 9,
        "十": 10,
        "十一": 11,
        "十二": 12,
        "十三": 13,
        "十四": 14,
        "十五": 15,
    }
    topic_map = [
        ("致辞", 2),
        ("ゲーム", 3),
        ("游戏", 3),
        ("广告", 5),
        ("広告", 5),
        ("旅游", 7),
        ("観光", 7),
        ("新聞", 9),
        ("新聞記事", 9),
        ("技術", 11),
        ("技术", 11),
        ("商务", 13),
        ("ビジネス", 13),
        ("信函", 13),
        ("文書", 13),
        ("法律", 15),
        ("法規", 15),
        ("招商", 15),
        ("招标", 15),
        ("招投标", 15),
        ("入札", 15),
        ("可能态", 9),
    ]
    if "第二学期" in text:
        for keyword, week in topic_map:
            if keyword in text:
                return week
    match = re.search(r"第(十一|十二|十三|十四|十五|一|二|三|四|五|六|七|八|九|十)[課课]", text)
    if match:
        return cn_map[match.group(1)]
    match = re.search(r"第(十一|十二|十三|十四|十五|一|二|三|四|五|六|七|八|九|十)章", text)
    if match:
        return cn_map[match.group(1)]
    match = re.search(r"第([0-9]{1,2})[課课]", text)
    if match:
        value = int(match.group(1))
        if 1 <= value <= 16:
            return value
    match = re.search(r"(?:^|/)([0-9]{1,2})[ ._、-]", text)
    if match:
        value = int(match.group(1))
        if 1 <= value <= 16:
            return value
    for keyword, week in topic_map:
        if keyword in text:
            return week
    return None


def summarize(text: str, title: str) -> str:
    lines = [line.strip() for line in text.splitlines() if len(line.strip()) >= 8]
    summary = "；".join(lines[:3])
    return (summary or title)[:260]


if __name__ == "__main__":
    main()
